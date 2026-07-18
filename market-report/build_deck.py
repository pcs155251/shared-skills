# -*- coding: utf-8 -*-
"""Build 預售屋分析 deck from a 實價登錄 pricelist xlsx.

Usage: uv run build_deck.py [pricelist.xlsx] [output.pptx]
(uv creates the venv from pyproject.toml on first run)
"""
import datetime
import os
import re
import sys
import tempfile
import pandas as pd
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm

SCRATCH = tempfile.mkdtemp(prefix="build_deck_")
XLSX = sys.argv[1] if len(sys.argv) > 1 else "PriceList (1).xlsx"
OUT = sys.argv[2] if len(sys.argv) > 2 else "大安區預售屋分析.pptx"

# =====================================================================
# MANUAL SUMMARY TEXT — the only qualitative judgment in this script.
# Regenerate with any AI (or edit by hand) whenever the dataset changes;
# {placeholders} are auto-filled from computed stats, so plain numbers
# never go stale — only the *interpretation* needs a fresh look.
# =====================================================================
MANUAL = dict(
    finding1_extra="銷售速度明顯領先同業。",
    finding3_extra="30-50坪中間帶(各約15%)亦相對疲弱。",
    summary4=(
        "{region_short}平均單價{avg_up:.0f}萬/坪、中位總價{med_total:,.0f}萬,"
        "「總價」才是去化關鍵——市場呈雙峰結構:低總價小宅(1-2房,鎖定置產與首購頂客)"
        "與50-70坪高階換屋(單價最高、達{maxbin_up:.0f}萬/坪)兩端熱絡,"
        "中間帶因總價已破8千萬但空間感有限而受壓縮。建議新案產品配比向2房與大坪數兩端傾斜,"
        "並以車位彈性搭售控制總價門檻。"),
    hist_headline2=(
        "總價呈雙峰:3,000-7,000萬小宅首購群 與 9,000萬-1.5億換屋豪宅群,"
        "中位數{med_total:,.0f}萬"),
)

# ---------- fonts (charts need a local CJK font; pick per platform) ----------
FONT_CANDIDATES = [  # (path, family name)
    ("/mnt/c/Windows/Fonts/msjh.ttc", "Microsoft JhengHei"),   # WSL
    ("/mnt/c/Windows/Fonts/msjhbd.ttc", "Microsoft JhengHei"),
    ("C:/Windows/Fonts/msjh.ttc", "Microsoft JhengHei"),       # Windows
    ("C:/Windows/Fonts/msjhbd.ttc", "Microsoft JhengHei"),
    ("/System/Library/Fonts/PingFang.ttc", "PingFang TC"),     # macOS
]
chart_font = None
for path, family in FONT_CANDIDATES:
    if os.path.exists(path):
        fm.fontManager.addfont(path)
        chart_font = chart_font or family
if chart_font is None:  # Linux fallback: any installed CJK-capable family
    installed = {f.name for f in fm.fontManager.ttflist}
    for family in ("Noto Sans TC", "Noto Sans CJK TC", "Noto Sans CJK JP",
                   "WenQuanYi Zen Hei", "AR PL UMing TW"):
        if family in installed:
            chart_font = family
            break
if chart_font is None:
    sys.exit("No CJK font found for charts — install Noto Sans TC "
             "(e.g. `sudo apt install fonts-noto-cjk` / `brew install font-noto-sans-cjk-tc`)")
plt.rcParams["font.family"] = chart_font
plt.rcParams["axes.unicode_minus"] = False

INK = "#0b0b0b"; INK2 = "#52514e"; MUTED = "#8a897f"
BLUE = "#2a78d6"; BLUE_L = "#ddeafc"; RED_L = "#fbe3dd"; RED = "#c0392b"
GRID = "#e8e7e3"

# ---------- data ----------
df = pd.read_excel(XLSX, header=1)
df.columns = [c.replace("\n", "") for c in df.columns]

def num(x):
    if pd.isna(x):
        return 0.0
    return float(str(x).replace(",", "").strip() or 0)

df["總價"] = df["總價(萬元)"].map(num)
df["車位價"] = df["車位總價(萬元)"].map(num)
df["單價"] = df["單價(萬元/坪)"].map(num)
df["權狀坪數"] = (df["總價"] - df["車位價"]) / df["單價"]
df["案名"] = df["建案名稱"].replace({"宏普建設台灣川普建設瑞閣": "宏普建設台灣川普瑞閣"})
df["解約"] = df["解約情形"].notna()

bins = [0, 20, 30, 40, 50, 70, 999]
bin_labels = ["20坪以下", "20-30坪", "30-40坪", "40-50坪", "50-70坪", "70坪以上"]
df["坪數區間"] = pd.cut(df["權狀坪數"], bins=bins, labels=bin_labels, right=False)

pv = (df.pivot_table(index="案名", columns="坪數區間", values="總價",
                     aggfunc="count", observed=True)
        .reindex(columns=bin_labels).fillna(0).astype(int))
pv["合計"] = pv.sum(axis=1)
pv = pv.sort_values("合計", ascending=False)

TOPN = 10
top_rows = pv.iloc[:TOPN]
other = pv.iloc[TOPN:].sum().to_frame().T
other.index = [f"其他({len(pv) - TOPN}案)"]
total = pv.sum().to_frame().T
total.index = ["合計"]
table_df = pd.concat([top_rows, other, total])

# ---------- computed stats for slide prose ----------
def fmt_price(x):  # 萬元 → '5,415萬' / '1.89億'
    return f"{x / 10000:.2f}億" if x >= 10000 else f"{x:,.0f}萬"

raw_hdr = str(pd.read_excel(XLSX, header=None, nrows=1).iloc[0, 0])
m = re.search(r"案件[::](.+?)(?=房地|建物|土地)", raw_hdr)
REGION = m.group(1).strip() if m else ""
region_short = REGION.replace("臺北市", "").replace("台北市", "") or REGION or "本區"
m = re.search(r"(\d{4})年(\d{2})月(\d{2})日", raw_hdr)
FETCH_DATE = f"{m.group(1)}/{m.group(2)}/{m.group(3)}" if m else "-"
TODAY = datetime.date.today().strftime("%Y/%m/%d")
IS_TAIPEI = "北市" in REGION

dates_all = df["交易日期"].astype(str).str.replace("　", "").str.strip()
PMIN, PMAX = dates_all.str[:6].min(), dates_all.str[:6].max()
n_all = len(df)
canc_all = int(df["解約"].sum())

top3 = list(pv["合計"].items())[:3] + [("-", 0)] * 3   # padded (name, count)
(p1, c1), (p2, c2), (p3, c3) = top3[0], top3[1], top3[2]
p1_canc = int(df.loc[df["案名"] == p1, "解約"].sum())
p1_dates = dates_all[df["案名"] == p1].str[:6]
p1_min, p1_max = p1_dates.min(), p1_dates.max()

bc = df["坪數區間"].value_counts()
b1, bc1 = bc.index[0], int(bc.iloc[0])
b2, bc2 = bc.index[1], int(bc.iloc[1])
b3, bc3 = bc.index[2], int(bc.iloc[2])
b1_median = df.loc[df["坪數區間"] == b1, "總價"].median()
b1_layout = df.loc[df["坪數區間"] == b1, "建物格局"].astype(str).str[:2].mode()[0]
bl = bc.index[-1]
bcl = int(bc.iloc[-1])
bl_sub = df[df["坪數區間"] == bl]
bl_median = bl_sub["總價"].median() if bcl else 0
bl_projs = "、".join(f"{k}{v}筆" for k, v in bl_sub["案名"].value_counts().head(2).items())

up_by_bin = df.groupby("坪數區間", observed=True)["單價"].mean()
STATS = dict(region_short=region_short,
             avg_up=df["單價"].mean(),
             med_total=df["總價"].median(),
             maxbin_up=up_by_bin.max())

# ---------- chart 1: 總價帶 histogram ----------
edges = [0, 3000, 4000, 5000, 6000, 7000, 8000, 9000, 10000,
         11000, 12000, 13000, 14000, 15000, 99999]
hlabels = (["3,000\n以下"]
           + [f"{a // 1000},000\n-{b // 1000},000" for a, b in zip(edges[1:-2], edges[2:-1])]
           + ["15,000\n以上"])
hcnt = pd.cut(df["總價"], bins=edges, right=False).value_counts().sort_index().values

n_all = len(df)
fig, ax = plt.subplots(figsize=(11.6, 4.6), dpi=200)
x = np.arange(len(hcnt))
ax.bar(x, hcnt, width=0.62, color=BLUE, zorder=3)
for xi, v in zip(x, hcnt):
    ax.text(xi, v + max(hcnt) * 0.115, str(v), ha="center", va="bottom",
            fontsize=11.5, color=INK)
    ax.text(xi, v + max(hcnt) * 0.02, f"{v / n_all * 100:.1f}%", ha="center",
            va="bottom", fontsize=8.5, color=INK2)
ax.set_xticks(x)
ax.set_xticklabels(hlabels, fontsize=9.5, color=INK2)
ax.set_ylabel("成交件數", fontsize=11, color=INK2)
ax.set_xlabel("總價帶(萬元)", fontsize=11, color=INK2)
ax.set_ylim(0, max(hcnt) * 1.28)
ax.yaxis.grid(True, color=GRID, zorder=0)
for s in ["top", "right", "left"]:
    ax.spines[s].set_visible(False)
ax.spines["bottom"].set_color(GRID)
ax.tick_params(colors=INK2, length=0)
fig.tight_layout()
fig.savefig(f"{SCRATCH}/hist.png", facecolor="white")
plt.close(fig)

# ---------- 消控表 target projects ----------
# argv[3:] = project names requested by the user; fall back to top-3 sellers.
proj_counts = df["案名"].value_counts()
requested = [a.strip() for a in sys.argv[3:] if a.strip()]
valid = [p for p in requested if p in proj_counts.index]
for p in requested:
    if p not in proj_counts.index:
        print(f"warning: 建案「{p}」不在pricelist內,略過")
grid_projects = valid or list(proj_counts.index[:3])
print("消控表建案:", ", ".join(grid_projects))


def parse_unit_floor(s):
    """'A6-14F號' / '15F-A3號' / 'A棟A5-12F號' / 'A3棟3F號' → (unit, floor)."""
    s = str(s).replace("號", "").replace("　", "").strip()
    m = re.search(r"(\d+)F", s)
    if not m:
        return None, None
    floor = int(m.group(1))
    unit = s.replace(m.group(0), "").replace("-", "").strip()
    if unit.endswith("棟"):          # 'A3棟3F號' → unit A3 (棟 suffix, no unit id)
        unit = unit[:-1]
    return unit or "?", floor

# ---------- pptx ----------
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

C_INK = RGBColor(0x0B, 0x0B, 0x0B)
C_INK2 = RGBColor(0x52, 0x51, 0x4E)
C_MUTED = RGBColor(0x8A, 0x89, 0x7F)
C_BLUE = RGBColor(0x2A, 0x78, 0xD6)
C_BLUE_D = RGBColor(0x1B, 0x4E, 0x8C)
C_HDRBG = RGBColor(0x1B, 0x4E, 0x8C)
C_ROW_A = RGBColor(0xFF, 0xFF, 0xFF)
C_ROW_B = RGBColor(0xF2, 0xF6, 0xFC)
C_TOTBG = RGBColor(0xDD, 0xEA, 0xFC)
FONT = "Microsoft JhengHei"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def add_text(slide, x, y, w, h, runs, size=18, bold=False, color=C_INK,
             align=PP_ALIGN.LEFT, line_spacing=1.15):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    if isinstance(runs, str):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        if isinstance(para, str):
            para = [(para, {})]
        for text, kw in para:
            r = p.add_run()
            r.text = text
            r.font.name = FONT
            r.font.size = Pt(kw.get("size", size))
            r.font.bold = kw.get("bold", bold)
            r.font.color.rgb = kw.get("color", color)
    return tb

def add_rule(slide, x, y, w, color=C_BLUE, h=0.045):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def slide_title(slide, title, subtitle=None):
    add_text(slide, 0.55, 0.32, 12.2, 0.7, title, size=27, bold=True)
    add_rule(slide, 0.58, 0.98, 1.5)
    if subtitle:
        add_text(slide, 0.58, 1.08, 12.1, 0.4, subtitle, size=13, color=C_MUTED)

# ----- slide 1: cover -----
s = prs.slides.add_slide(BLANK)
add_rule(s, 0.9, 2.35, 2.2, C_BLUE, 0.07)
add_text(s, 0.85, 2.6, 11.6, 1.0, f"{REGION} 預售屋市場分析".strip(), size=40, bold=True)
add_text(s, 0.9, 3.7, 11.6, 0.5, "坪數結構 × 總價帶分布 × 熱銷建案消控",
         size=20, color=C_BLUE_D, bold=True)
add_text(s, 0.9, 5.9, 11.6, 0.9, [
    f"資料來源:不動產交易實價查詢服務網(預售屋,{PMIN}-{PMAX},共{n_all}筆)",
    f"資料擷取:{FETCH_DATE} 製表:{TODAY}"], size=12.5, color=C_MUTED)

# ----- slide 2: 坪數區間 table -----
s = prs.slides.add_slide(BLANK)
slide_title(s, "坪數區間 × 建案 成交統計",
            f"房屋權狀坪數 =(總價 − 車位總價)÷ 單價;統計期間 {PMIN}-{PMAX},含解約{canc_all}筆")
rows, cols = len(table_df) + 1, 8
tw, th = 12.23, 5.5
tbl_shape = s.shapes.add_table(rows, cols, Inches(0.55), Inches(1.55), Inches(tw), Inches(th))
tbl = tbl_shape.table
tbl.columns[0].width = Inches(3.13)
for c in range(1, 8):
    tbl.columns[c].width = Inches(1.3)

def set_cell(cell, text, *, bold=False, color=C_INK, size=12, align=PP_ALIGN.CENTER, fill=None):
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = cell.margin_right = Emu(45720)
    cell.margin_top = cell.margin_bottom = Emu(9144)
    p = cell.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color

headers = ["建案名稱"] + bin_labels + ["合計"]
for c, htxt in enumerate(headers):
    set_cell(tbl.cell(0, c), htxt, bold=True, color=RGBColor(255, 255, 255),
             size=12.5, fill=C_HDRBG, align=PP_ALIGN.CENTER if c else PP_ALIGN.LEFT)
for ri, (name, row) in enumerate(table_df.iterrows(), start=1):
    is_total = name == "合計"
    fill = C_TOTBG if is_total else (C_ROW_A if ri % 2 else C_ROW_B)
    set_cell(tbl.cell(ri, 0), str(name), bold=is_total, size=12, fill=fill, align=PP_ALIGN.LEFT)
    vals = [row[b] for b in bin_labels] + [row["合計"]]
    mx = max(row[b] for b in bin_labels)
    for ci, v in enumerate(vals, start=1):
        hot = (not is_total) and ci <= 6 and v == mx and v > 0
        set_cell(tbl.cell(ri, ci), str(int(v)) if v else "-",
                 bold=is_total or hot or ci == 7,
                 color=C_BLUE_D if hot else (C_INK if v else C_MUTED),
                 size=12, fill=fill)
add_text(s, 0.55, 7.08, 12.2, 0.35,
         "藍色粗體=該建案最熱銷坪數帶"
         + (f";「其他」彙總{len(pv) - TOPN}個成交{int(pv['合計'].iloc[TOPN])}筆以下之建案"
            if len(pv) > TOPN else ""),
         size=10.5, color=C_MUTED)

# ----- slide 3: 總結 -----
s = prs.slides.add_slide(BLANK)
slide_title(s, "坪數區間分析 總結")
findings = [
    ("1", f"賣最好的建案:{p1}",
     f"統計期間內成交{c1}筆(含{p1_canc}筆解約,淨{c1 - p1_canc}筆),"
     f"佔全區{c1 / n_all * 100:.1f}%,居冠;其次為{p2}({c2}筆)、{p3}({c3}筆)。"
     f"{p1}自{p1_min}至{p1_max}即去化{c1}戶," + MANUAL["finding1_extra"]),
    ("2", f"賣最多的坪數帶:{b1}",
     f"共{bc1}筆、佔{bc1 / n_all * 100:.1f}%,為市場絕對主力;"
     f"其次為{b2}({bc2}筆,{bc2 / n_all * 100:.1f}%)與{b3}({bc3}筆,{bc3 / n_all * 100:.1f}%)。"
     f"中位總價{fmt_price(b1_median)}的{b1_layout}產品是{region_short}預售最大宗。"),
    ("3", f"最難賣的坪數帶:{bl}",
     f"僅{bcl}筆、佔{bcl / n_all * 100:.1f}%({bl_projs}),"
     f"中位總價高達{fmt_price(bl_median)},客群極窄、去化最慢;" + MANUAL["finding3_extra"]),
    ("4", "市場總結",
     MANUAL["summary4"].format(**STATS)),
]
y = 1.5
for tag, head, body in findings:
    add_rule(s, 0.62, y + 0.1, 0.045, C_BLUE, 1.08 if tag != "4" else 1.45)
    add_text(s, 0.85, y, 11.9, 0.4,
             [[(f"{tag}  ", {"color": C_BLUE, "bold": True}), (head, {"bold": True})]], size=16.5)
    add_text(s, 1.12, y + 0.42, 11.6, 0.8, body, size=12.5, color=C_INK2, line_spacing=1.25)
    y += 1.32 if tag != "4" else 1.7

# ----- slide 4: histogram -----
hi = int(np.argmax(hcnt))
mode_label = hlabels[hi].replace("\n", "")
mode_label = mode_label if mode_label.endswith("下") or mode_label.endswith("上") \
    else mode_label + "萬"
share_low = (df["總價"] < 8000).mean() * 100
s = prs.slides.add_slide(BLANK)
slide_title(s, "市場總價帶分布(臺北市:3,000萬以下起,每1,000萬一區間)")
add_text(s, 0.58, 1.28, 12.2, 0.4,
         [[(f"主力總價帶為{mode_label}({hcnt[hi]}筆、{hcnt[hi] / n_all * 100:.1f}%),"
            f"{share_low:.0f}%的成交落在8,000萬以下",
            {"bold": True, "color": C_BLUE_D})]], size=16)
add_text(s, 0.58, 1.72, 12.2, 0.4,
         [[(MANUAL["hist_headline2"].format(**STATS),
            {"bold": True, "color": C_BLUE_D})]], size=16)
s.shapes.add_picture(f"{SCRATCH}/hist.png", Inches(0.75), Inches(2.35), width=Inches(11.85))
add_text(s, 0.58, 7.1, 12.2, 0.3,
         f"n={n_all}(含解約{canc_all}筆)"
         + (";本區全屬臺北市,故一律採臺北市分區規則" if IS_TAIPEI else ""),
         size=10.5, color=C_MUTED)

# ----- slides 5+: 消控表, one per requested project (native editable tables) -----
C_SOLD = RGBColor(0xDD, 0xEA, 0xFC)
C_RECENT = RGBColor(0xD6, 0xF0, 0xE0)
C_CANC = RGBColor(0xFB, 0xE3, 0xDD)
C_EMPTY = RGBColor(0xF7, 0xF7, 0xF5)
C_RED = RGBColor(0xC0, 0x39, 0x2B)
C_GREEN = RGBColor(0x1B, 0x6E, 0x45)

def set_grid_cell(cell, lines, *, fill, colors, sizes, bold=False):
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = cell.margin_right = Emu(18288)   # 0.02"
    cell.margin_top = cell.margin_bottom = Emu(0)
    tf = cell.text_frame
    tf.word_wrap = False
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = 0.9
        r = p.add_run()
        r.text = line
        r.font.name = FONT
        r.font.size = Pt(sizes[i] if isinstance(sizes, (list, tuple)) else sizes)
        r.font.bold = bold
        r.font.color.rgb = colors[i] if isinstance(colors, list) else colors

for rank, proj in enumerate(grid_projects, start=1):
    f = df[df["案名"] == proj].copy()
    parsed = f["棟及號"].map(parse_unit_floor)
    f["戶別"] = parsed.str[0]
    f["樓層"] = parsed.str[1]
    f = f[f["樓層"].notna()].copy()
    f["樓層"] = f["樓層"].astype(int)
    if f.empty:
        print(f"warning: 「{proj}」的棟及號無法解析樓層,略過消控表")
        continue
    units = sorted(f["戶別"].unique())
    floors = list(range(int(f["樓層"].max()), int(f["樓層"].min()) - 1, -1))
    f["_date"] = f["交易日期"].astype(str).str.replace("　", "").str.strip()
    f["年月"] = f["_date"].str[:6]
    latest_ym = f["年月"].max()
    # same unit can appear twice (解約後再售) — let the non-解約, newest record win
    f = f.sort_values(["解約", "_date"], ascending=[False, True])
    cell = {(r["戶別"], r["樓層"]): r for _, r in f.iterrows()}

    hdr = {}
    for u in units:
        ping = f.loc[f["戶別"] == u, "權狀坪數"]
        layout = f.loc[f["戶別"] == u, "建物格局"].astype(str).str[:2].mode()[0]
        hdr[u] = f"{u}戶({layout} 約{ping.min():.0f}-{ping.max():.0f}坪)" \
            if ping.max() - ping.min() > 2 else f"{u}戶({layout} 約{ping.mean():.0f}坪)"

    # adaptive sizing: many units → narrower columns, many floors → one-line cells
    one_line = len(floors) > 22
    fs1, fs2 = (7.5, 6.5) if len(units) <= 6 else (6.5, 6.0)
    if one_line:
        fs1 = fs2 = min(fs2, 6.0)

    s = prs.slides.add_slide(BLANK)
    rank_all = list(proj_counts.index).index(proj) + 1
    slide_title(s, f"{proj} 銷售消控表(成交{len(f)}筆,全區第{rank_all}名)",
                f"橫軸=戶別(房型),縱軸=樓層(揭露範圍 {floors[-1]}F-{floors[0]}F);"
                "格內:總價 / 單價 / 成交日期")

    nrows, ncols = len(floors) + 1, len(units) + 1
    gt = s.shapes.add_table(nrows, ncols, Inches(0.4), Inches(1.5),
                            Inches(12.53), Inches(5.7)).table
    gt.first_row = False
    gt.horz_banding = False
    gt.columns[0].width = Inches(0.53)
    for c in range(1, ncols):
        gt.columns[c].width = Inches(12.0 / len(units))
    gt.rows[0].height = Inches(0.3)
    for r in range(1, nrows):
        gt.rows[r].height = Inches(min(0.32, 5.4 / len(floors)))

    set_grid_cell(gt.cell(0, 0), ["樓層"], fill=C_HDRBG,
                  colors=RGBColor(255, 255, 255), sizes=9, bold=True)
    for j, u in enumerate(units, start=1):
        set_grid_cell(gt.cell(0, j), [hdr[u]], fill=C_HDRBG,
                      colors=RGBColor(255, 255, 255),
                      sizes=9.5 if len(units) <= 6 else 8, bold=True)

    for i, fl in enumerate(floors, start=1):
        set_grid_cell(gt.cell(i, 0), [f"{fl}F"], fill=RGBColor(0xEE, 0xF2, 0xF8),
                      colors=C_INK2, sizes=8.5 if not one_line else 6.5, bold=True)
        for j, u in enumerate(units, start=1):
            r = cell.get((u, fl))
            if r is None:
                set_grid_cell(gt.cell(i, j), [""], fill=C_EMPTY, colors=C_MUTED, sizes=7)
                continue
            d = str(r["交易日期"]).replace("　", "").strip()
            canc = bool(r["解約"])
            recent = (not canc) and r["年月"] == latest_ym
            l1 = f"{r['總價']:,.0f}萬  {r['單價']:.1f}萬/坪"
            l2 = d + "(已解約)" if canc else d
            if canc:
                fill, cols = C_CANC, [C_RED, C_RED]
            elif recent:
                fill, cols = C_RECENT, [C_GREEN, C_GREEN]
            else:
                fill, cols = C_SOLD, [C_INK, C_INK2]
            if one_line:
                set_grid_cell(gt.cell(i, j), [f"{l1}  {l2}"], fill=fill,
                              colors=[cols[0]], sizes=[fs1])
            else:
                set_grid_cell(gt.cell(i, j), [l1, l2], fill=fill, colors=cols,
                              sizes=[fs1, fs2])

    recs = list(cell.values())  # count what the grid shows, not raw deal rows
    n_disclosed = len(recs)
    n_canc = sum(bool(r["解約"]) for r in recs)
    n_recent = sum(1 for r in recs if not r["解約"] and r["年月"] == latest_ym)
    n_sold = n_disclosed - n_recent - n_canc
    n_total = len(units) * len(floors)
    n_unsold = n_total - n_disclosed
    add_text(s, 0.4, 7.22, 12.5, 0.3, [[
        (f"已成交(實價登錄揭露){n_sold}戶", {"color": C_BLUE_D, "bold": True}),
        ("   ", {}),
        (f"最近一個月({latest_ym})成交 {n_recent}戶", {"color": C_GREEN, "bold": True}),
        ("   ", {}),
        (f"已解約 {n_canc}戶", {"color": C_RED, "bold": True}),
        ("   ", {}),
        (f"未揭露(未售或尚未登錄){n_unsold}戶", {"color": C_MUTED, "bold": True}),
        (f"  |  以{len(units)}戶×{len(floors)}層({floors[-1]}-{floors[0]}F)"
         f"推估{n_total}戶,已揭露{n_disclosed}戶,去化約{n_disclosed / n_total * 100:.0f}%",
         {"color": C_MUTED}),
    ]], size=10)

prs.save(OUT)
print("saved", OUT)
print("table rows:", len(table_df))
